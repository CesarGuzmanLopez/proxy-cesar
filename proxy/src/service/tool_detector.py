"""Content handling utilities for unsupported content types.

When the user sends content (images, audio, files) to a model that lacks
the required capability, the proxy stores the base64 data as a blob in
Valkey and replaces it with a text reference containing size, type, and
a brief description so the model knows what the user sent.

Images from the same message are described together in a single vision
model call, using the user's original text prompt for context.
"""

import asyncio
import base64
import hashlib
import json
import logging
import re

import fitz  # PyMuPDF — mandatory for PDF text extraction

from src.config.constants import BLOB_STORAGE_TTL_SECONDS


# ── Lenient JSON parser (handles trailing commas from Llama 4 Scout, etc.) ──

_JSON_TRAILING_COMMA_RE = re.compile(r",\s*([}\]])")
_JSON_CTRL_CHAR_RE = re.compile(r"[\000-\037]")


def _parse_json_lenient(text: str):
    """Parse JSON with lenient handling for common LLM output issues.

    - Strips trailing commas before ``]`` and ``}`` (common in Llama 4 Scout output)
    - Strips unescaped control characters (newlines, tabs) inside strings
    - Falls back to strict ``json.loads`` if cleaned version also fails
    """
    # 1. Strip trailing commas
    cleaned = _JSON_TRAILING_COMMA_RE.sub(r"\1", text)
    # 2. Strip unescaped control characters inside strings
    cleaned = _JSON_CTRL_CHAR_RE.sub("", cleaned)
    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        return json.loads(text)  # Fallback: maybe original was fine

logger = logging.getLogger(__name__)

BLOB_PREFIX = "BLOB"
BLOB_TTL = BLOB_STORAGE_TTL_SECONDS  # Configurable in src/config/constants.py
_MAX_BLOB_SIZE = 10 * 1024 * 1024  # 10 MB
_BLOB_FORMAT_VERSION = "v2"  # Increment when format changes


# ── Content Type Classification (SINGLE source of truth) ────────────


class ContentType:
    """Content types recognized by the proxy's extraction pipeline."""
    IMAGE = "image"
    AUDIO = "audio"
    PDF = "pdf"
    DOCUMENT = "document"  # Word, text, CSV, etc.
    SPREADSHEET = "spreadsheet"  # Excel, etc.
    PRESENTATION = "presentation"  # PowerPoint
    UNKNOWN = "unknown"


def _classify_content_type(part: dict) -> str:
    """Classify a content part into a ContentType. SINGLE source of truth.

    Handles all formats:
    - OpenAI:  {"type": "image_url", "image_url": {"url": "..."}}
    - File:    {"type": "file", "file": {"mime_type": "...", "data": "..."}}
    - Anthropic: {"type": "file", "source": {"media_type": "...", "data": "..."}}
    - Audio:   {"type": "input_audio", ...}
    - Text with data URL: {"type": "text", "text": "data:image/..."}
    """
    ptype = part.get("type", "")
    if ptype in ("image_url", "image"):
        return ContentType.IMAGE
    if ptype == "input_audio":
        return ContentType.AUDIO
    if ptype == "text":
        text = part.get("text", "")
        if isinstance(text, str) and text.startswith("data:image/"):
            return ContentType.IMAGE
        return ContentType.UNKNOWN
    if ptype != "file":
        return ContentType.UNKNOWN

    # File type — determine from MIME
    mime = (
        part.get("mime_type")
        or part.get("mimeType")
        or part.get("mimetype")
        or part.get("media_type", "")
    )
    if not mime:
        file_obj = part.get("file", {}) or {}
        mime = (
            file_obj.get("mime_type")
            or file_obj.get("mimeType")
            or file_obj.get("mimetype")
            or file_obj.get("media_type", "")
        )
    if not mime:
        source = part.get("source", {}) or {}
        mime = source.get("media_type", "")
    if not mime:
        # Try extracting from data URI
        for field in ("data", "file", "source"):
            candidate = part.get(field)
            if isinstance(candidate, str) and candidate.startswith("data:"):
                match = re.match(r"data:([a-z]+/[a-z0-9+.-]+)", candidate)
                if match:
                    mime = match.group(1)
                    break
            elif isinstance(candidate, dict):
                inner = candidate.get("data") or candidate.get("file_data", "")
                if isinstance(inner, str) and inner.startswith("data:"):
                    match = re.match(r"data:([a-z]+/[a-z0-9+.-]+)", inner)
                    if match:
                        mime = match.group(1)
                        break

    mime_lower = mime.lower() if mime else ""
    if "pdf" in mime_lower:
        return ContentType.PDF
    if any(v in mime_lower for v in ("spreadsheet", "excel", "xls", "xlsx", "csv", "ods")):
        return ContentType.SPREADSHEET
    if any(v in mime_lower for v in ("presentation", "powerpoint", "ppt", "pptx")):
        return ContentType.PRESENTATION
    if any(v in mime_lower for v in ("word", "docx", "officedocument", "opendocument", "text/", "rtf", "msword")):
        return ContentType.DOCUMENT
    if any(v in mime_lower for v in ("image", "png", "jpg", "jpeg", "gif", "webp")):
        return ContentType.IMAGE
    if any(v in mime_lower for v in ("audio", "wav", "mp3", "ogg", "flac", "m4a")):
        return ContentType.AUDIO
    if any(v in mime_lower for v in ("video", "mp4", "webm", "mkv", "avi")):
        return ContentType.UNKNOWN  # Not supported for extraction

    return ContentType.DOCUMENT  # Fallback to document for unknown file types


def _hash_content(data: str) -> str:
    return hashlib.sha256(data.encode()).hexdigest()[:16]


def _extract_mime(data_uri: str) -> str | None:
    match = re.match(r"data:([a-z]+/[a-z0-9+.-]+)", data_uri)
    return match.group(1) if match else None


def _find_model_with_capability(config, cap: str = "vision"):
    """Find any configured physical model with a given capability.

    Returns the full ``PhysicalModelSchema`` object (not just the model string)
    so callers can access ``api_base`` and ``api_key_env`` for custom endpoints.

    TODO: pick the cheapest, not the first. Cost data in pseudo_models.yaml notes.
    """
    for pm in config.pseudo_models.values():
        for phys in pm.physical_models:
            if getattr(phys, cap, False):
                return phys
    return None


def _extract_user_text(content: list[dict]) -> str:
    """Join all text parts from a user message for context."""
    return " ".join(
        str(p.get("text", ""))
        for p in content
        if isinstance(p, dict) and p.get("type") == "text"
    )


def _resolve_api_key(phys) -> str | None:
    """Resolve API key from environment if the physical model has api_key_env set."""
    if not phys or not phys.api_key_env:
        return None
    import os

    return os.environ.get(phys.api_key_env) or None


async def _describe_image_batch(
    images: list[tuple[str, str]], user_prompt: str, config
) -> list[str]:
    """Describe multiple images together in one vision model call.

    Degrades images before sending (reduces quality to save tokens).
    Checks context window capacity and splits into batches if needed.
    Returns one description per image. Returns empty strings on failure.

    If a single image exceeds the model's context capacity, returns an
    error description so the user knows the images are too large.
    """
    vision_phys = _find_model_with_capability(config, "vision")
    if not vision_phys:
        return [""] * len(images)
    if not images:
        return []

    from src.service.multimedia.image_processor import (
        can_batch_fit,
        degrade_image_async,
        estimate_image_tokens,
    )

    vision_model = vision_phys.model
    context_window = vision_phys.context_window or 128000
    # Degrade all images first (async — non-blocking)
    degraded: list[tuple[str, str]] = []
    for h, raw in images:
        degraded_raw = await degrade_image_async(raw)
        degraded.append((h, degraded_raw))

    system = (
        "You receive images that the user sent to a model without vision. "
        "For EACH image, analyze exhaustively for a text-only AI model:\n"
        "1. Extract ALL visible text exactly as it appears (including crossed out, "
        "strikethrough, handwritten, annotations, headers, footers, labels, watermarks). "
        "Preserve exact spelling and original language.\n"
        "2. Describe the image in rich detail: scene, objects, people, actions, "
        "expressions, clothing, environment, lighting, colors, composition, style, "
        "spatial relationships, perspective.\n"
        "3. For UI/code/diagrams: describe layout, structure, hierarchy, "
        "interactive elements, data values, axes, legends, trends.\n"
        "4. Note imperfections: blur, glare, artifacts, damage, low resolution areas.\n"
        "5. If multiple panels/sections, describe each separately.\n\n"
        "Return a JSON array of strings in image order, one full analysis per image. "
        'Format per element: "[EXTRACTED TEXT]: ...\\n[DESCRIPTION]: ...\\n[TECHNICAL DETAILS]: ..."'
    )

    # Estimate if all images fit in one batch
    sample_tokens = estimate_image_tokens(degraded[0][1], "high")
    text_tokens = len(user_prompt) // 4 + 200  # rough text estimate

    if not can_batch_fit(len(degraded), sample_tokens, context_window, text_tokens):
        # Split into smaller batches
        from src.service.multimedia.image_processor import _MAX_IMAGES_PER_BATCH

        results: list[str] = []
        for i in range(0, len(degraded), _MAX_IMAGES_PER_BATCH):
            batch = degraded[i : i + _MAX_IMAGES_PER_BATCH]
            batch_tokens = estimate_image_tokens(batch[0][1], "high")
            if not can_batch_fit(len(batch), batch_tokens, context_window, text_tokens):
                # Single batch still too large — return error descriptions
                logger.warning(
                    "batch_too_large model=%s images=%d ctx=%d",
                    vision_model,
                    len(batch),
                    context_window,
                )
                results.extend(
                    [
                        "[ERROR: Image too large to process. Reduce size or use a vision-capable model directly.]"
                    ]
                    * len(batch)
                )
                continue
            batch_descs = await _describe_single_batch(
                batch,
                user_prompt,
                system,
                vision_model,
                vision_phys,
            )
            results.extend(batch_descs)
        return results

    return await _describe_single_batch(
        degraded,
        user_prompt,
        system,
        vision_model,
        vision_phys,
    )


async def _describe_single_batch(
    images: list[tuple[str, str]],
    user_prompt: str,
    system: str,
    vision_model: str,
    vision_phys,
) -> list[str]:
    """Send a single batch of images to the vision model for description."""
    content: list[dict] = [
        {
            "type": "text",
            "text": f"User: {user_prompt}\nDescribe each image below in JSON array format:",
        }
    ]
    for _, raw in images:
        content.append({"type": "image_url", "image_url": {"url": raw}})

    try:
        from src.adapters.litellm import call_litellm

        response = await call_litellm(
            model=vision_model,
            messages=[
                {"role": "system", "content": system},
                {"role": "user", "content": content},
            ],
            api_base=vision_phys.api_base or None,
            api_key=_resolve_api_key(vision_phys),
            max_tokens=2048 * len(images),
            temperature=0.1,
        )
        resp = response.model_dump() if hasattr(response, "model_dump") else response
        if isinstance(resp, dict):
            text = resp.get("choices", [{}])[0].get("message", {}).get("content", "")
            if text:
                text = text.strip()
                if text.startswith("```"):
                    text = text.split("\n", 1)[-1].rsplit("```", 1)[0].strip()
                parsed = _parse_json_lenient(text)
                if isinstance(parsed, list):
                    return [str(d) for d in parsed]
    except Exception as exc:
        logger.warning(
            "batch_describe_failed model=%s images=%d error=%s",
            vision_model,
            len(images),
            str(exc),
        )
    return [""] * len(images)


async def _transcribe_audio(raw_data: str, config) -> str:
    """Transcribe audio using Whisper via litellm.atranscription()."""
    audio_phys = _find_model_with_capability(config, "audio")
    if not audio_phys or not raw_data.startswith("data:"):
        return ""
    audio_model = audio_phys.model
    try:
        audio_bytes = base64.b64decode(raw_data.split(",", 1)[-1])
        from litellm import atranscription
        from io import BytesIO

        audio_file = BytesIO(audio_bytes)
        audio_file.name = "audio.wav"
        response = await atranscription(
            model=audio_model,
            file=audio_file,
            api_base=audio_phys.api_base or None,
            api_key=_resolve_api_key(audio_phys),
            temperature=0.1,
        )
        return (getattr(response, "text", None) or "")[:500]
    except Exception as exc:
        logger.warning(
            "audio_transcribe_failed model=%s error=%s", audio_model, str(exc)
        )
        return ""


async def _try_extract_pdf_text(base64_data: str) -> str:
    """Extract text from a base64-encoded PDF using PyMuPDF (offloaded to thread pool).

    Returns a JSON structure with per-page details:
      - Page number and content type (text / blank / image_only / mixed)
      - Text preview (up to 2000 chars per page)
      - Embedded image count
      - Per-page notes (blank, image-only, truncation)
    Plus document-level statistics (total pages, blank pages, image-only pages, etc.).

    The JSON format allows the LLM to understand document structure, identify
    which pages contain relevant information, and detect blank or image-only pages.
    """

    _PER_PAGE_TEXT_MAX = 4000
    _MAX_DETAILED_PAGES = 100

    def _sync_extract() -> str:
        try:
            pdf_bytes = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            page_count = len(doc)
            size_kb = len(pdf_bytes) // 1024

            pages: list[dict] = []
            pages_with_text = 0
            pages_blank = 0
            pages_image_only = 0
            total_chars = 0

            for i, page in enumerate(doc):
                if i >= _MAX_DETAILED_PAGES:
                    break
                page_num = i  # 0-indexed (PyMuPDF convention)
                text = page.get_text().strip()
                images = page.get_images()
                has_images = len(images) > 0
                char_count = len(text)
                total_chars += char_count

                if not text and not has_images:
                    page_type = "blank"
                    pages_blank += 1
                elif not text and has_images:
                    page_type = "image_only"
                    pages_image_only += 1
                elif text and not has_images:
                    page_type = "text"
                    pages_with_text += 1
                else:
                    page_type = "mixed"
                    pages_with_text += 1

                entry: dict = {"page": page_num, "type": page_type}

                if text:
                    if char_count <= _PER_PAGE_TEXT_MAX:
                        entry["text_preview"] = text
                    else:
                        entry["text_preview"] = text[:_PER_PAGE_TEXT_MAX]
                        entry["note"] = (
                            f"Text truncated to first {_PER_PAGE_TEXT_MAX} chars "
                            f"(page has {char_count} total)"
                        )
                if has_images:
                    entry["image_count"] = len(images)

                if page_type == "blank":
                    entry["note"] = "No visible text or images on this page"
                elif page_type == "image_only":
                    entry["note"] = (
                        f"Contains {len(images)} embedded image(s) but no "
                        "extractable text; may require OCR"
                    )

                pages.append(entry)

            doc.close()

            remaining_pages = page_count - min(page_count, _MAX_DETAILED_PAGES)

            result = {
                "document_type": "PDF",
                "total_pages": page_count,
                "size_kb": size_kb,
                "extraction_method": "PyMuPDF",
                "statistics": {
                    "pages_with_text": pages_with_text,
                    "blank_pages": pages_blank,
                    "image_only_pages": pages_image_only,
                    "total_extracted_chars": total_chars,
                    "detailed_pages_shown": len(pages),
                    "pages_not_shown": remaining_pages,
                },
                "pages": pages,
            }
            return json.dumps(result, indent=2, ensure_ascii=False)
        except Exception as exc:
            return json.dumps({
                "error": f"PDF could not parse: {str(exc)[:200]}",
            })

    return await asyncio.to_thread(_sync_extract)


async def _try_extract_docx_text(base64_data: str) -> str:
    """Extract text from a base64-encoded DOCX file (offloaded to thread pool).

    Returns a JSON structure with document-level metadata, paragraph/table
    counts, and text preview grouped into logical segments so the model can
    understand the document structure and decide whether to use its own tools.
    """

    _DOCX_PREVIEW_CHARS = 4000
    _MAX_TABLE_ROWS_PREVIEW = 30

    def _sync_extract() -> str:
        try:
            from docx import Document
            from io import BytesIO

            docx_bytes = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            doc = Document(BytesIO(docx_bytes))

            text_parts = [para.text for para in doc.paragraphs if para.text.strip()]
            total_chars = sum(len(p) for p in text_parts)
            para_with_text = len(text_parts)
            table_count = len(doc.tables)
            size_kb = len(docx_bytes) // 1024

            # Group paragraphs into segments of ~2000 chars
            segments: list[dict] = []
            seg_buf: list[str] = []
            seg_chars = 0
            seg_idx = 0
            for p in text_parts:
                seg_buf.append(p)
                seg_chars += len(p)
                if seg_chars >= 2000:
                    seg_idx += 1
                    seg_text = "\n".join(seg_buf)
                    segments.append({"segment": seg_idx, "paragraphs": len(seg_buf), "chars": len(seg_text), "text_preview": seg_text[:1500]})
                    seg_buf = []
                    seg_chars = 0
            if seg_buf:
                seg_idx += 1
                seg_text = "\n".join(seg_buf)
                segments.append({"segment": seg_idx, "paragraphs": len(seg_buf), "chars": len(seg_text), "text_preview": seg_text[:1500]})

            total_paragraphs = len(doc.paragraphs)
            segments_shown = len(segments)

            # Extract table previews
            tables_preview: list[dict] = []
            for ti, table in enumerate(doc.tables[:5]):  # max 5 tables in preview
                rows: list[str] = []
                for ri, row in enumerate(table.rows):
                    if ri >= _MAX_TABLE_ROWS_PREVIEW:
                        break
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                tables_preview.append({
                    "table": ti + 1,
                    "rows": len(table.rows),
                    "cols": len(table.columns),
                    "preview_rows": rows[:10],
                })

            result = {
                "document_type": "DOCX",
                "total_paragraphs": total_paragraphs,
                "paragraphs_with_text": para_with_text,
                "tables": table_count,
                "size_kb": size_kb,
                "extraction_method": "python-docx",
                "statistics": {
                    "total_extracted_chars": total_chars,
                    "segments_shown": segments_shown,
                    "tables_previewed": len(tables_preview),
                },
                "segments": segments,
                "tables_preview": tables_preview,
            }
            return json.dumps(result, indent=2, ensure_ascii=False)
        except ImportError:
            return json.dumps({"error": "DOCX: python-docx library not available"})
        except Exception as exc:
            return json.dumps({"error": f"DOCX could not parse: {str(exc)[:200]}"})

    return await asyncio.to_thread(_sync_extract)


async def _try_extract_pptx_text(base64_data: str) -> str:
    """Extract text from a base64-encoded PPTX file (offloaded to thread pool).

    Returns a JSON structure with per-slide detail so the model can understand
    the presentation structure, identify which slides have content, and decide
    whether to use its own tools.
    """

    _SLIDE_TEXT_MAX = 4000
    _MAX_DETAILED_SLIDES = 100

    def _sync_extract() -> str:
        try:
            from pptx import Presentation
            from io import BytesIO

            pptx_bytes = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            prs = Presentation(BytesIO(pptx_bytes))
            size_kb = len(pptx_bytes) // 1024
            total_slides = len(prs.slides)

            slides: list[dict] = []
            slides_with_text = 0
            slides_empty = 0
            total_chars = 0

            for i, slide in enumerate(prs.slides):
                if i >= _MAX_DETAILED_SLIDES:
                    break
                slide_num = i + 1
                slide_parts: list[str] = []
                has_table = False
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_parts.append(text)
                    if shape.has_table:
                        has_table = True
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            slide_parts.append(" | ".join(row_text))
                slide_text = "\n".join(slide_parts).strip()
                char_count = len(slide_text)
                total_chars += char_count

                if not slide_text:
                    slides_empty += 1
                    slides.append({
                        "slide": slide_num,
                        "has_text": False,
                        "note": "Empty slide (no extractable text)",
                    })
                else:
                    slides_with_text += 1
                    entry: dict = {"slide": slide_num, "has_text": True}
                    if char_count <= _SLIDE_TEXT_MAX:
                        entry["text_preview"] = slide_text
                    else:
                        entry["text_preview"] = slide_text[:_SLIDE_TEXT_MAX]
                        entry["note"] = (
                            f"Text truncated to first {_SLIDE_TEXT_MAX} chars "
                            f"(slide has {char_count} total)"
                        )
                    if has_table:
                        entry["has_table"] = True
                    slides.append(entry)

            remaining = total_slides - min(total_slides, _MAX_DETAILED_SLIDES)
            result = {
                "document_type": "PPTX",
                "total_slides": total_slides,
                "size_kb": size_kb,
                "extraction_method": "python-pptx",
                "statistics": {
                    "slides_with_text": slides_with_text,
                    "empty_slides": slides_empty,
                    "total_extracted_chars": total_chars,
                    "detailed_slides_shown": len(slides),
                    "slides_not_shown": remaining,
                },
                "slides": slides,
            }
            return json.dumps(result, indent=2, ensure_ascii=False)
        except ImportError:
            return json.dumps({"error": "PPTX: python-pptx library not available"})
        except Exception as exc:
            return json.dumps({"error": f"PPTX could not parse: {str(exc)[:200]}"})

    return await asyncio.to_thread(_sync_extract)


async def _try_extract_xlsx_text(base64_data: str) -> str:
    """Extract text from a base64-encoded XLSX file (offloaded to thread pool).

    Returns a JSON structure with per-sheet detail so the model can understand
    the spreadsheet structure, identify which sheets have data, and decide
    whether to use its own tools.
    """

    _ROWS_PER_SHEET_PREVIEW = 100
    _MAX_DETAILED_SHEETS = 50

    def _sync_extract() -> str:
        try:
            import openpyxl
            from io import BytesIO

            xlsx_bytes = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            wb = openpyxl.load_workbook(BytesIO(xlsx_bytes), read_only=True, data_only=True)
            size_kb = len(xlsx_bytes) // 1024
            total_sheets = len(wb.sheetnames)
            total_data_rows = 0

            sheets: list[dict] = []
            for si, sheet_name in enumerate(wb.sheetnames):
                if si >= _MAX_DETAILED_SHEETS:
                    break
                ws = wb[sheet_name]
                row_values_list: list[str] = []
                sheet_row_count = 0
                for row in ws.iter_rows(values_only=True):
                    row_values = [str(v) if v is not None else "" for v in row]
                    row_text = " | ".join(row_values).strip()
                    if row_text:
                        row_values_list.append(row_text)
                        sheet_row_count += 1
                        total_data_rows += 1

                if sheet_row_count > _ROWS_PER_SHEET_PREVIEW:
                    preview = row_values_list[:_ROWS_PER_SHEET_PREVIEW]
                    sheets.append({
                        "sheet": sheet_name,
                        "data_rows": sheet_row_count,
                        "preview_rows_shown": _ROWS_PER_SHEET_PREVIEW,
                        "text_preview": "\n".join(preview),
                        "note": f"Rows truncated to first {_ROWS_PER_SHEET_PREVIEW} (sheet has {sheet_row_count} total)",
                    })
                else:
                    sheets.append({
                        "sheet": sheet_name,
                        "data_rows": sheet_row_count,
                        "preview_rows_shown": sheet_row_count,
                        "text_preview": "\n".join(row_values_list) if row_values_list else "",
                    })

            wb.close()
            remaining = total_sheets - min(total_sheets, _MAX_DETAILED_SHEETS)

            result = {
                "document_type": "XLSX",
                "total_sheets": total_sheets,
                "size_kb": size_kb,
                "extraction_method": "openpyxl",
                "statistics": {
                    "total_data_rows": total_data_rows,
                    "detailed_sheets_shown": len(sheets),
                    "sheets_not_shown": remaining,
                },
                "sheets": sheets,
            }
            return json.dumps(result, indent=2, ensure_ascii=False)
        except ImportError:
            return json.dumps({"error": "XLSX: openpyxl library not available"})
        except Exception as exc:
            return json.dumps({"error": f"XLSX could not parse: {str(exc)[:200]}"})

    return await asyncio.to_thread(_sync_extract)


def _classify_content_parts(  # noqa: S3776 — multiple content types × multiple formats
    content: list,
) -> tuple[
    str,
    list[tuple[str, str, str, str, str]],
    list[tuple[str, str, str, str, str]],
    list[tuple[str, str, str, str, str]],
    list[dict],
]:
    """Classify content parts into images, audio, files, and others.

    Supports multiple image formats:
    - OpenAI: {"type": "image_url", "image_url": {"url": "..."}}
    - Base64: {"type": "image", "image": "data:image/..."}
    - Text with data URL: {"type": "text", "text": "data:image/..."}
    """
    user_text = _extract_user_text(content)
    images: list[tuple[str, str, str, str, str]] = []
    audios: list[tuple[str, str, str, str, str]] = []
    files: list[tuple[str, str, str, str, str]] = []
    others: list[dict] = []

    for part in content:
        if not isinstance(part, dict):
            others.append(part)
            continue

        ptype = part.get("type", "")
        raw = ""
        content_type = None

        # Extract raw data from various formats
        if ptype == "image_url":
            # OpenAI standard format
            raw = part.get("image_url", {}).get("url", "")
            content_type = "image"
        elif ptype == "image":
            # Base64 inline format (used by Anthropic, others)
            raw = part.get("image", "")
            content_type = "image"
        elif ptype == "text":
            # Some clients send images as text with data: URL
            text_content = part.get("text", "")
            if isinstance(text_content, str) and text_content.startswith("data:image/"):
                raw = text_content
                content_type = "image"
            else:
                # Regular text, not an image
                others.append(part)
                continue
        elif ptype == "input_audio":
            raw = (part.get("input_audio", {}) or {}).get("data", "")
            content_type = "audio"
        elif ptype == "file":
            # Support multiple file formats:
            # OpenAI: {"type": "file", "file": {"data": "data:...", "mime_type": "..."}}
            # Anthropic: {"type": "file", "source": {"type": "base64", "media_type": "...", "data": "..."}}
            # Direct: {"type": "file", "data": "data:...", "mimeType": "..."}
            raw = (
                (part.get("file", {}) or {}).get("data")
                or (part.get("file", {}) or {}).get("file_data")
                or (part.get("source", {}) or {}).get("data", "")
                or part.get("data", "")
            )
            content_type = "file"
            # If no data found but source has base64 format, reconstruct
            if not raw and "source" in part:
                src = part["source"]
                if isinstance(src, dict):
                    raw = f"data:{src.get('media_type', 'application/octet-stream')};base64,{src.get('data', '')}"
        else:
            others.append(part)
            continue

        # Validate and process extracted data
        if not raw or not raw.startswith("data:"):
            others.append(part)
            continue

        h = _hash_content(raw)
        mime = _extract_mime(raw) or f"{content_type}/unknown"
        sz = str(len(raw) // 1024)
        filename = ""
        info = (h, raw, mime, sz, filename)

        # Classify into appropriate category
        if content_type == "image":
            images.append(info)
        elif content_type == "audio":
            audios.append(info)
        elif content_type == "file":
            files.append(info)

    return user_text, images, audios, files, others


def _truncate_desc(desc: str, sz_kb: int | str) -> str:
    """Truncate description so the full blob text ≤ original file size in tokens."""
    if not desc:
        return desc
    sz = int(sz_kb) if isinstance(sz_kb, str) else sz_kb
    max_desc_chars = max(sz * 1024 - 400, 500)
    if len(desc) > max_desc_chars:
        return desc[:max_desc_chars] + "..."
    return desc


def _log_extraction_event(
    label: str,
    filename: str,
    sz: int | str,
    extraction_method: str,
    desc: str,
) -> None:
    """Log extraction event at INFO level with the extracted content."""
    sz_str = f"{sz} KB" if isinstance(sz, (int, float)) else str(sz)
    preview = desc[:300] if desc else "(empty)"
    logger.info(
        "extraction_event label=%s filename=%s size=%s method=%s "
        "desc_len=%d preview=%s",
        label,
        filename or "(unnamed)",
        sz_str,
        extraction_method or "unknown",
        len(desc) if desc else 0,
        preview,
    )


def _build_blob_text(label: str, sz: int | str, desc: str = "", extraction_method: str = "", filename: str = "") -> str:
    """Build extracted content text block.

    Args:
        label: Content type label (image, document, pdf, etc.)
        sz: File size in KB
        desc: Extracted text description/content
        extraction_method: How the content was extracted
        filename: Original filename if available

    Format is versioned with _BLOB_FORMAT_VERSION for future compatibility.
    """
    desc = _truncate_desc(desc, sz)
    t = f"[{_BLOB_FORMAT_VERSION}][File extracted: {label}"
    if filename:
        t += f" | name: {filename}"
    t += f" | {sz} KB"
    if extraction_method:
        t += f" | tool: {extraction_method}"
    t += "]\n"
    t += "Sent as file — read with tools if in workspace, or use extracted content below:\n"
    if desc:
        t += f"\n{desc}\n"
    else:
        t += "\n(Warning: Content extraction failed.)\n"

    # Log the final blob text that the LLM will receive
    logger.info(
        "extraction_blob label=%s filename=%s size=%s method=%s blob_len=%d",
        label,
        filename or "(unnamed)",
        f"{sz} KB" if isinstance(sz, (int, float)) else str(sz),
        extraction_method or "unknown",
        len(t),
    )
    return t


_EXTRACTION_LABELS = {
    "image": "Vision model (Llama 4 Scout / MiMo Omni)",
    "audio": "Speech-to-text (Whisper)",
}


def _determine_doc_extraction(mime: str) -> str:
    """Determine extraction method label based on MIME type."""
    mime_lower = mime.lower()
    if "pdf" in mime_lower:
        return "PyMuPDF (Python)"
    if any(v in mime_lower for v in ("wordprocessingml", "word", "docx", "msword", "opendocument")):
        return "python-docx (Python)"
    if any(v in mime_lower for v in ("spreadsheet", "excel", "xls", "xlsx", "csv", "ods")):
        return "openpyxl (Python)"
    if any(v in mime_lower for v in ("presentation", "powerpoint", "ppt", "pptx")):
        return "python-pptx (Python)"
    return "text decode (UTF-8)"


def _build_blob_output(others, images, descs, audios, aresults, files, fresults):
    """Build output content list from classified parts and descriptions."""
    out: list[dict[str, object]] = list(others)

    for (_, _, mime, sz, filename), d in zip(images, descs):
        text = _build_blob_text("image", sz, d, _EXTRACTION_LABELS["image"], filename)
        out.append({"type": "text", "text": text})

    for (_, _, mime, sz, filename), d in zip(audios, aresults):
        text = _build_blob_text("audio", sz, d, _EXTRACTION_LABELS["audio"], filename)
        out.append({"type": "text", "text": text})

    for (_, _, mime, sz, filename), d in zip(files, fresults):
        extraction = _determine_doc_extraction(mime)
        text = _build_blob_text("document", sz, d, extraction, filename)
        out.append({"type": "text", "text": text})

    return out


async def _process_msg_blobs(
    msg: dict[str, object],
    prefix: str,
    valkey,
    config,
) -> list[dict[str, object]]:
    """Process one user message: classify, store, describe, build output."""
    content = msg.get("content", "")
    if not isinstance(content, list):
        return [msg]

    # Pass 1: classify and store
    user_text, image_blobs, audio_blobs, file_blobs, other_parts = (
        _classify_content_parts(content)
    )
    store_tasks = [
        _store_blob_if_missing(valkey, f"{prefix}:{h}", r)
        for h, r, _, _, _ in image_blobs + audio_blobs + file_blobs
    ]
    if store_tasks:
        await asyncio.gather(*store_tasks)

    # Pass 2: describe
    descriptions = await _describe_images(
        valkey, prefix, image_blobs, user_text, config
    )
    audio_results = (
        await asyncio.gather(
            *[
                _describe_audio(valkey, f"{prefix}:{h}:desc", r, config)
                for h, r, _, _, _ in audio_blobs
            ]
        )
        if audio_blobs
        else []
    )
    pdf_results = (
        await asyncio.gather(
            *[
                _describe_file_generic(valkey, f"{prefix}:{h}:desc", r, mime)
                for h, r, mime, _, _ in file_blobs
            ]
        )
        if file_blobs
        else []
    )

    # Pass 3: build output
    out = _build_blob_output(
        other_parts,
        image_blobs,
        descriptions,
        audio_blobs,
        audio_results,
        file_blobs,
        pdf_results,
    )
    return [{**msg, "content": out}]


async def _describe_images(valkey, prefix, blobs, user_text, config):
    """Batch describe images or return cached descriptions.

    Cache key includes a hash of the user prompt so that the same image
    asked with a different question gets a different (correct) description.
    Key format: {prefix}:{hash}:desc:{prompt_hash}

    NOTE: Image descriptions are NOT deterministic because they depend on
    the user's prompt context — same image + different prompt → new description.
    """
    if not blobs:
        return []
    prompt_hash = hashlib.sha256((user_text or "").encode()).hexdigest()[:8]
    cache_keys = [f"{prefix}:{h}:desc:{prompt_hash}" for h, _, _, _, _ in blobs]
    cached = await asyncio.gather(
        *[_get_cached(valkey, k) for k in cache_keys]
    )

    # Log cache results per image
    all_cached = True
    for (h, _, _, _, fn), k, v in zip(blobs, cache_keys, cached):
        if v:
            logger.info(
                "extraction_image cache_hit=true key=%s filename=%s "
                "desc_len=%d preview=%s",
                k, fn or "(unnamed)",
                len(v), v[:200] if v else "(empty)",
            )
        else:
            all_cached = False
            logger.info(
                "extraction_image cache_hit=false key=%s filename=%s "
                "prompt_hash=%s",
                k, fn or "(unnamed)", prompt_hash,
            )

    if all_cached:
        return cached

    descs = await _describe_image_batch(
        [(h, r) for h, r, _, _, _ in blobs], user_text, config
    )
    while len(descs) < len(blobs):
        descs.append("")

    # Log newly described images
    for (h, _, _, _, fn), d in zip(blobs, descs):
        if d:
            logger.info(
                "extraction_image cache_hit=false key=%s filename=%s "
                "desc_len=%d prompt_hash=%s preview=%s",
                f"{prefix}:{h}:desc:{prompt_hash}",
                fn or "(unnamed)",
                len(d), prompt_hash,
                d[:300] if d else "(empty)",
            )

    store = [
        _store_desc(valkey, f"{prefix}:{h}:desc:{prompt_hash}", d)
        for (h, _, _, _, _), d in zip(blobs, descs)
        if d
    ]
    if store:
        await asyncio.gather(*store)
    return descs


async def _get_cached(valkey, key: str) -> str:
    try:
        v = await valkey.get(key)
        return v or ""
    except Exception as exc:
        logger.debug("blob_cache_get_error key=%s err=%s", key, exc)
        return ""


async def replace_base64_with_blob_refs(
    messages: list[dict[str, object]],
    conversation_id: str | None = None,
    valkey=None,
    config=None,
) -> list[dict[str, object]]:
    """Replace base64 content with [BLOB:hash:mime] references.

    Groups images from the same message and describes them together
    using the user's original prompt for context.
    Real URLs pass through unchanged.
    """
    if valkey is None:
        return messages
    prefix = f"blob:{conversation_id or 'anon'}"
    results: list[dict[str, object]] = []
    for msg in messages:
        if msg.get("role") != "user":
            results.append(msg)
        else:
            results.extend(await _process_msg_blobs(msg, prefix, valkey, config))
    return results


async def _store_blob_if_missing(valkey, key: str, raw: str) -> None:
    """Store blob only if it doesn't exist yet."""
    try:
        exists = await valkey.exists(key)
        if not exists:
            await valkey.set(key, raw, ex=BLOB_TTL)
            logger.info(
                "cache_write type=blob key=%s size=%d ttl=%d",
                key, len(raw), BLOB_TTL,
            )
        else:
            logger.debug("cache_skip type=blob key=%s reason=already_exists", key)
    except Exception as exc:
        logger.warning("cache_write_error type=blob key=%s err=%s", key, exc)


async def _store_desc(valkey, key: str, desc: str) -> None:
    """Store description with SETNX + EXPIRE."""
    try:
        stored = await valkey.setnx(key, desc)
        await valkey.expire(key, BLOB_TTL)
        if stored:
            logger.info(
                "cache_write type=desc key=%s desc_len=%d ttl=%d",
                key, len(desc), BLOB_TTL,
            )
        else:
            logger.debug("cache_skip type=desc key=%s reason=already_exists", key)
    except Exception as exc:
        logger.warning("cache_write_error type=desc key=%s err=%s", key, exc)


async def _describe_audio(valkey, desc_key: str, raw: str, config) -> str:
    """Transcribe audio if not already cached (SHA256 hash → deterministic)."""
    try:
        cached = await valkey.get(desc_key)
        if cached:
            logger.info(
                "extraction_audio cache_hit=true desc_key=%s desc_len=%d preview=%s",
                desc_key,
                len(cached),
                cached[:300] if cached else "(empty)",
            )
            return cached
    except Exception as exc:
        logger.warning("blob_audio_cache_error key=%s err=%s", desc_key, exc)
    desc = await _transcribe_audio(raw, config)
    logger.info(
        "extraction_audio cache_hit=false desc_key=%s desc_len=%d preview=%s",
        desc_key,
        len(desc) if desc else 0,
        desc[:300] if desc else "(empty)",
    )
    if desc:
        await _store_desc(valkey, desc_key, desc)
    return desc


async def _describe_file_generic(valkey, desc_key: str, raw: str, mime: str) -> str:
    """Extract text from a file — PDF, DOCX, PPTX, XLSX, or plain text fallback.

    Cached by SHA256 hash of raw content (deterministic extraction).
    """
    try:
        cached = await valkey.get(desc_key)
        if cached:
            logger.info(
                "extraction_file cache_hit=true desc_key=%s mime=%s desc_len=%d preview=%s",
                desc_key,
                mime,
                len(cached),
                cached[:400] if cached else "(empty)",
            )
            return cached
    except Exception as exc:
        logger.warning("blob_file_cache_error key=%s err=%s", desc_key, exc)

    # Use a mock part dict with mime_type and data to classify
    mock_part = {"type": "file", "file": {"mime_type": mime, "data": raw}}
    ctype = _classify_content_type(mock_part)

    if ctype == ContentType.PDF:
        desc = await _try_extract_pdf_text(raw)
    elif ctype == ContentType.DOCUMENT:
        desc = await _try_extract_docx_text(raw)
    elif ctype == ContentType.PRESENTATION:
        desc = await _try_extract_pptx_text(raw)
    elif ctype == ContentType.SPREADSHEET:
        desc = await _try_extract_xlsx_text(raw)
    else:
        # Generic text fallback: decode base64 as UTF-8
        try:
            decoded = base64.b64decode(raw.split(",", 1)[-1].strip())
            text = decoded.decode("utf-8", errors="replace")[:5000]
            desc = f"[Text extracted ({len(decoded)} bytes)]\n\n{text}" if text.strip() else ""
        except Exception as exc:
            logger.debug("blob_text_decode_error mime=%s err=%s", mime, exc)
            desc = ""

    logger.info(
        "extraction_file cache_hit=false desc_key=%s mime=%s ctype=%s desc_len=%d preview=%s",
        desc_key,
        mime,
        ctype,
        len(desc) if desc else 0,
        desc[:500] if desc else "(empty)",
    )
    if desc:
        await _store_desc(valkey, desc_key, desc)
    return desc


def inject_blob_extraction_guidance(messages: list[dict]) -> list[dict]:
    """Inject system message explaining blob extraction and tool availability.

    When the proxy auto-extracts content (describes images, transcribes audio,
    extracts PDF text), it sends the extracted information in the message.
    This system message tells the model:

    1. Content was auto-extracted from blobs (multimodal content the model can't process)
    2. The model can choose to use its own specialized tools for alternative analysis
    3. How to access the blob reference if needed for custom processing

    This allows models with specialized tools (PDF parsers, image recognition,
    audio analysis) to decide whether to use the proxy's extraction or invoke
    their own tools for custom or more accurate analysis.

    Returns:
        messages with injected system message if blobs detected, otherwise unchanged
    """

    def _message_has_blobs(msg: dict) -> bool:
        content = msg.get("content", "")
        if isinstance(content, str):
            return "[v2][File extracted:" in content
        if isinstance(content, list):
            return any(
                isinstance(part, dict)
                and isinstance(part.get("text"), str)
                and ("[v2][File extracted:" in part["text"])
                for part in content
            )
        return False

    if not any(_message_has_blobs(m) for m in messages):
        return messages

    if any(m.get("role") == "system" for m in messages):
        return messages

    # Detect which content types are in the messages
    counts: dict[str, int] = {}
    for msg in messages or []:
        content = msg.get("content", [])
        if isinstance(content, str):
            text = content.lower()
            if "[v2][file extracted: image" in text:
                counts["images"] = counts.get("images", 0) + 1
            if any(t in text for t in ("[v2][file extracted: pdf", "[v2][file extracted: document", "[v2][file extracted: presentation", "[v2][file extracted: spreadsheet")):
                counts["docs"] = counts.get("docs", 0) + 1
            if "[v2][file extracted: audio" in text:
                counts["audios"] = counts.get("audios", 0) + 1
        elif isinstance(content, list):
            for item in content:
                if not isinstance(item, dict):
                    continue
                ctype = _classify_content_type(item)
                if ctype == ContentType.IMAGE:
                    counts["images"] = counts.get("images", 0) + 1
                elif ctype == ContentType.AUDIO:
                    counts["audios"] = counts.get("audios", 0) + 1
                elif ctype in (ContentType.PDF, ContentType.DOCUMENT, ContentType.PRESENTATION, ContentType.SPREADSHEET):
                    counts["docs"] = counts.get("docs", 0) + 1
                # Also check formatted text for already-processed messages
                if item.get("type") == "text":
                    text = str(item.get("text", "")).lower()
                    if "[v2][file extracted: image" in text:
                        counts["images"] = counts.get("images", 0) + 1
                    if any(t in text for t in ("[v2][file extracted: pdf", "[v2][file extracted: document", "[v2][file extracted: presentation", "[v2][file extracted: spreadsheet")):
                        counts["docs"] = counts.get("docs", 0) + 1
                    if "[v2][file extracted: audio" in text:
                        counts["audios"] = counts.get("audios", 0) + 1

    methods = []
    if counts.get("images", 0) > 0:
        methods.append("  • Images → Vision models (Llama 4 Scout / MiMo Omni)")
    if counts.get("audios", 0) > 0:
        methods.append("  • Audio → Whisper speech-to-text")
    if counts.get("docs", 0) > 0:
        methods.append("  • Documents → PyMuPDF / python-docx / openpyxl / python-pptx (Python)")

    methods_text = "\n".join(methods) if methods else (
        "  • Images → Vision models (Llama 4 Scout / MiMo Omni)\n"
        "  • Audio → Whisper speech-to-text\n"
        "  • Documents → PyMuPDF / python-docx / openpyxl / python-pptx (Python)"
    )

    system_message = (
        "**File Content Extraction**\n\n"
        "Some files were auto-extracted. Their content is in [v2][File extracted: ...] blocks below. "
        "Read directly, or access original files with tools if they exist in your workspace.\n\n"
        f"Extraction methods:\n{methods_text}"
    )

    return [{"role": "system", "content": system_message}] + messages

