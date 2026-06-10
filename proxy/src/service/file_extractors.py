"""Document/text file extraction utilities for base64-encoded content.

Extracts text from PDF, DOCX, PPTX, and XLSX files using dedicated
Python libraries (PyMuPDF, python-docx, python-pptx, openpyxl).
All I/O-heavy operations are offloaded to a thread pool.

No circular dependency: this module imports ONLY from stdlib and
3rd-party extraction libraries — never from tool_detector or other
first-party modules.
"""

import asyncio
import base64
import json
import re

import fitz  # PyMuPDF — mandatory for PDF text extraction


async def _try_extract_pdf_text(base64_data: str) -> str:
    """Extract text from a base64-encoded PDF using PyMuPDF (offloaded to thread pool).

    Returns a JSON structure with per-page details:
      - Page number and content type (text / blank / image_only / mixed)
      - Text preview (up to 2000 chars per page)
      - Embedded image count
      - Per-page notes (blank, image-only, truncation)
    Plus document-level statistics (total pages, blank pages, image-only pages, etc.).

    The JSON format allows the LLM to understand document structure, identify
    which pages contain relevant information, and detect blank or image-only pages.
    """

    _PER_PAGE_TEXT_MAX = 4000
    _MAX_DETAILED_PAGES = 100

    def _sync_extract() -> str:
        try:
            pdf_bytes = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            page_count = len(doc)
            size_kb = len(pdf_bytes) // 1024

            pages: list[dict] = []
            pages_with_text = 0
            pages_blank = 0
            pages_image_only = 0
            total_chars = 0

            for i in range(min(page_count, _MAX_DETAILED_PAGES)):
                page = doc[i]
                text = page.get_text()
                images = page.get_images(full=True)
                image_count = len(images)
                clean = " ".join(text.split())[:_PER_PAGE_TEXT_MAX]
                page_info: dict = {"page": i + 1}
                if clean:
                    page_info["type"] = "text"
                    page_info["content"] = clean
                    pages_with_text += 1
                elif image_count:
                    page_info["type"] = "image_only"
                    page_info["page_content_images"] = image_count
                    pages_image_only += 1
                else:
                    page_info["type"] = "blank"
                    notes = []
                    if image_count > 0:
                        notes.append(f"{image_count} images")
                    pages_blank += 1

                page_text_size = len(text.strip())
                if page_text_size > 0:
                    total_chars += page_text_size
                    if not clean and not image_count and not page_text_size:
                        page_info["notes"] = notes if notes else ["blank"]
                if image_count:
                    page_info["page_images"] = image_count
                pages.append(page_info)

            has_more_pages = page_count > _MAX_DETAILED_PAGES
            stats = {
                "document_type": "PDF",
                "extraction_method": "PyMuPDF",
                "size_kb": size_kb,
                "total_pages": page_count,
                "pages_with_text": pages_with_text,
                "blank_pages": pages_blank,
                "image_only_pages": pages_image_only,
                "total_chars_extracted": total_chars,
                "truncated": has_more_pages,
            }
            if has_more_pages:
                stats["note"] = f"Only showing details for first {_MAX_DETAILED_PAGES} of {page_count} pages"

            return json.dumps(stats | {"detailed_pages": pages}, ensure_ascii=False)
        except Exception as exc:
            return json.dumps({
                "document_type": "PDF",
                "extraction_method": "PyMuPDF",
                "error": f"PDF could not parse: {str(exc)[:200]}",
            })

    return await asyncio.to_thread(_sync_extract)


async def _try_extract_docx_text(base64_data: str) -> str:
    """Extract text from a base64-encoded DOCX file using python-docx."""
    from docx import Document as DocxDocument

    def _sync_extract() -> str:
        try:
            raw = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            doc = DocxDocument(raw)
            size_kb = len(raw) // 1024

            # Extract all paragraphs
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Extract all tables
            tables = []
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                tables.append("\n".join(rows))

            text = "\n".join(paragraphs)
            if tables:
                text += "\n\n--- TABLES ---\n" + "\n\n".join(tables)

            return json.dumps({
                "document_type": "Word Document",
                "extraction_method": "python-docx",
                "size_kb": size_kb,
                "paragraphs": len(paragraphs),
                "tables": len(tables),
                "content": text[:50000],
                "truncated": len(text) > 50000,
            }, ensure_ascii=False)
        except Exception as exc:
            return json.dumps({
                "document_type": "Word Document",
                "extraction_method": "python-docx",
                "error": str(exc)[:200],
            })

    return await asyncio.to_thread(_sync_extract)


async def _try_extract_pptx_text(base64_data: str) -> str:
    """Extract text from a base64-encoded PPTX file using python-pptx."""
    from pptx import Presentation

    def _sync_extract() -> str:
        try:
            raw = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            prs = Presentation(raw)
            size_kb = len(raw) // 1024
            slides_text: list[str] = []

            for i, slide in enumerate(prs.slides, 1):
                slide_parts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = para.text.strip()
                            if txt:
                                slide_parts.append(txt)
                    if shape.has_table:
                        table = shape.table
                        rows = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows.append(" | ".join(cells))
                        slide_parts.append("TABLE:\n" + "\n".join(rows))
                if slide_parts:
                    slides_text.append(f"--- Slide {i} ---\n" + "\n".join(slide_parts))

            full = "\n\n".join(slides_text)
            return json.dumps({
                "document_type": "PowerPoint Presentation",
                "extraction_method": "python-pptx",
                "size_kb": size_kb,
                "slides": len(slides_text),
                "content": full[:30000],
                "truncated": len(full) > 30000,
            }, ensure_ascii=False)
        except Exception as exc:
            return json.dumps({
                "document_type": "PowerPoint Presentation",
                "extraction_method": "python-pptx",
                "error": str(exc)[:200],
            })

    return await asyncio.to_thread(_sync_extract)


async def _try_extract_xlsx_text(base64_data: str) -> str:
    """Extract text from a base64-encoded XLSX file using openpyxl."""
    from openpyxl import load_workbook

    def _sync_extract() -> str:
        try:
            raw = base64.b64decode(base64_data.split(",", 1)[-1].strip())
            wb = load_workbook(raw, data_only=True, read_only=True)
            size_kb = len(raw) // 1024
            sheets_text: list[str] = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    rows.append(" | ".join(cells))
                if rows:
                    sheets_text.append(f"--- Sheet: {sheet_name} ({len(ws[1]) if ws.max_column else 0} cols, {len(rows)} rows) ---\n" + "\n".join(rows))
            wb.close()

            full = "\n\n".join(sheets_text)
            return json.dumps({
                "document_type": "Excel Spreadsheet",
                "extraction_method": "openpyxl",
                "size_kb": size_kb,
                "sheets": len(sheets_text),
                "content": full[:30000],
                "truncated": len(full) > 30000,
            }, ensure_ascii=False)
        except Exception as exc:
            return json.dumps({
                "document_type": "Excel Spreadsheet",
                "extraction_method": "openpyxl",
                "error": str(exc)[:200],
            })

    return await asyncio.to_thread(_sync_extract)
